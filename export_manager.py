from reportlab.lib.pagesizes import landscape, LETTER
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import os

class ExportManager:
    @staticmethod
    def export_as_pdf(scenes, svg_paths, output_path):
        c = canvas.Canvas(output_path, pagesize=landscape(LETTER))
        width, height = landscape(LETTER)

        for i, (scene, svg_path) in enumerate(zip(scenes, svg_paths)):
            # Draw SVG
            try:
                drawing = svg2rlg(svg_path)
                # Scale drawing to fit
                factor = (width - 2*inch) / drawing.width
                drawing.width *= factor
                drawing.height *= factor
                drawing.scale(factor, factor)
                
                renderPDF.draw(drawing, c, 1*inch, height - drawing.height - 1*inch)
            except Exception as e:
                c.drawString(1*inch, height - 2*inch, f"Error rendering SVG: {str(e)}")

            # Draw Metadata
            c.setFont("Helvetica-Bold", 14)
            c.drawString(1*inch, 2*inch, f"SCENE {i+1}: {scene.header}")
            
            c.setFont("Helvetica", 10)
            c.drawString(1*inch, 1.7*inch, f"SHOT: {scene.shot_type} | MOVEMENT: {scene.camera_movement}")
            c.drawString(1*inch, 1.5*inch, f"LIGHTING: {scene.lighting}")
            
            # Dialogue/Action
            textobject = c.beginText(1*inch, 1.2*inch)
            textobject.setFont("Courier", 9)
            textobject.textLines(scene.dialogue[:200] + "..." if len(scene.dialogue) > 200 else scene.dialogue)
            c.drawText(textobject)

            c.showPage()
        c.save()

    @staticmethod
    def export_as_pptx(scenes, png_paths, output_path):
        prs = Presentation()
        # Set to 16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        for i, (scene, png_path) in enumerate(zip(scenes, png_paths)):
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
            
            # Use the PNG preview since PPT doesn't support SVG natively without complex conversion
            try:
                if os.path.exists(png_path):
                    slide.shapes.add_picture(png_path, Inches(1), Inches(0.5), width=Inches(11.33))
            except Exception as e:
                print(f"PPT Export error: {e}")
                slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(4)).text = f"[Image Error]"

            # Footer info
            txBox = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(2.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            p = tf.add_paragraph()
            p.text = f"SCENE {i+1}: {scene.header}"
            p.font.bold = True
            p.font.size = Pt(24)
            
            p2 = tf.add_paragraph()
            p2.text = f"{scene.shot_type} | {scene.camera_movement} | {scene.lighting}"
            p2.font.size = Pt(14)
            
            if scene.dialogue:
                p3 = tf.add_paragraph()
                p3.text = f"DIALOGUE:\n{scene.dialogue}"
                p3.font.size = Pt(12)
                p3.font.italic = True

        prs.save(output_path)
